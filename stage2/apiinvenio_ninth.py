import os
import time
import logging
import requests
import fitz
import pytesseract
from PIL import Image
from docx import Document
#from doc2docx import convert
import pptx
import sys # Deal with large text field in the output csv
import csv
import stat
from tika import parser
# Note: Apache Tika is what I turned to after failing many times dealing with .doc files.

# Set up tika locally
# Download the server `tika-server-standard-3.0.0.jar` from tika's website: https://tika.apache.org/download.html
os.environ['TIKA_SERVER_ENDPOINT'] = 'http://localhost:9998'

# Make sure large text field in csv can be processed
csv.field_size_limit(sys.maxsize)

# Set up logging and make it display in my console
logging.basicConfig(filename='process.log', level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
console_handler = logging.StreamHandler()
console_handler.setLevel(logging.INFO)
formatter = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
console_handler.setFormatter(formatter)
logging.getLogger().addHandler(console_handler)

def download_file(url, local_filename, auth_headers, max_retries=5, backoff_factor=1):
    """
    Downloads a file from the specified URL and saves it locally.

    Args:
        url (str): The URL of the file to download.
        local_filename (str): The path where the file will be saved.
        auth_headers (dict): HTTP headers containing authentication information.
        max_retries (int, optional): Maximum number of retry attempts. 5 is the default value.
        backoff_factor (int, optional): Delay factor for retries. 1 is the default value.
    
    Returns:
        str: The path to the downloaded file, or None if the download fails.
    """
    for attempt in range(max_retries):
        try:
            logging.info(f"Attempting to download {url} (Attempt {attempt + 1})")
            response = requests.get(url, headers=auth_headers, allow_redirects=True) # Send HTTP request to the specific url
            response.raise_for_status() # Check for HTTP errors
            
            # Check for redirect URL in the response
            if "http" in response.text:
                file_url = response.text.strip()
                logging.info(f"Fetching file content from URL: {file_url}")
                file_response = requests.get(file_url, allow_redirects=True)
                file_response.raise_for_status()
                
                with open(local_filename, 'wb') as file: # Open local_filename in write mode
                    for chunk in file_response.iter_content(chunk_size=8192): # Reads the response body in chunks of 8 KB
                        file.write(chunk)

            else:
                with open(local_filename, 'wb') as file:
                    for chunk in response.iter_content(chunk_size=8192):
                        file.write(chunk)

            logging.info(f"Successfully downloaded {local_filename}")
            return local_filename # Exits the function download_file

        except requests.exceptions.RequestException as e:
            logging.error(f"Error on attempt {attempt + 1}: {e}")
            time.sleep(backoff_factor * (2 ** attempt)) # Delay increases with each failed attempt (1, 2, 4, 8s)

    logging.error(f"Failed to download {url} after {max_retries} attempts")
    return None
 
def extract_file(local_file_path, file_name):
    """
    Extract text from a file based on its extension.

    Args:
        local_file_path (str): The path to the file.
        file_name (str): The name of the file.
    
    Returns:
        Extracted text (or None) and Flag (0 for success, 1 for failure).
    """
    # Switched to dictionary{} instead of list[]
    # List iterates through each entry, costs more time
    extraction_methods = {
        '.pdf': extract_text_from_pdf,
        '.docx': extract_text_from_docx,
        #'.doc': extract_text_from_doc, old way/had issues
        '.pptx': extract_text_from_pptx,
        '.txt': extract_text_from_txt,
        '.doc': extract_text_from_doc_with_tika,
    }

    # Extract file extension
    file_extension = os.path.splitext(file_name)[1].lower()
    logging.info(f"File extension detected: {file_extension}")

    # # Handle .doc files separately for conversion
    # if file_extension == '.doc':
    #     logging.info(f"Attempting to convert .doc file {file_name} to .docx.")
    #     newdocx_file_path = convert_doc_to_docx(local_file_path)
        
    #     # If conversion fails, log and return
    #     if not newdocx_file_path:
    #         logging.warning(f"Failed to convert {file_name} to .docx. Retained in folder for manual inspection.")
    #         return None, 1

    #     # Update file path and extension for further processing
    #     local_file_path = newdocx_file_path
    #     file_extension = '.docx'

    # Find the matching function for the file extension
    extraction_function = extraction_methods.get(file_extension)

    if extraction_function:
        try:
            logging.info(f"Extracting text from {file_name} using {extraction_function.__name__}") # Return the name of the function as a string
            extracted_text = extraction_function(local_file_path)
            if extracted_text:
                return extracted_text, 0 # Success #TODO: Return 3 values, extracted_text, failed, supported; add third value of 1 (supported)
            else:
                return None, 1 # Failure #TODO: third value should be 1 (supported)
        except Exception as e:
            logging.error(f"Error extracting text from {file_name} at {local_file_path}: {e}")
            return None, 1 #TODO: make changes, 1 (third value, supported)
    else:
        # Handle unsupported files
        logging.warning(f"Unsupported file type for {file_name}. Skipping this file.")
        return None, 1, #TODO:0(meaning not supported)
    
def extract_text_from_pdf(pdf_file_path):
    """
    Extract text from a PDF file.
    Uses PyMuPDF/fitz for text-based PDFs and PyTesseract for image-based PDFs.

    Args:
        pdf_file_path (str): The path to the PDF file.

    Returns:
        str: The extracted text, or None if it fails.
    """
    try:
        with fitz.open(pdf_file_path) as doc:
            text = ""
            total_pages = len(doc)
            for page_num in range(total_pages): # Iterates through each page
                if (page_num + 1) % 10 == 0 or page_num == total_pages - 1: # Log every 10 pages or the last page
                    logging.info(f"Processing page {page_num + 1} of {total_pages}")

                page = doc.load_page(page_num) # Load the current page into memory

                try:
                    page_text = page.get_text()
                    if page_text.strip(): #If there is text on the page, it's likely text-based
                        logging.debug(f"Page {page_num + 1}: Text-based")
                        text += page_text
                    else: #This is probably image-based
                        logging.debug(f"Page {page_num + 1}: Image-based")
                        pix = page.get_pixmap()
                        img = Image.frombytes("RGB",[pix.width, pix.height], pix.samples)
                        ocr_text = pytesseract.image_to_string(img)

                        # Validate OCR results: remove whitespace, measure length
                        if len(ocr_text.strip()) < 10: #TODO: thinking about if 10 is a good choice
                            logging.warning(f"OCR results on page {page_num + 1} may be incomplete.")
                            ocr_text = f"[Potential OCR issue on page {page_num + 1}]"
                        
                        text += ocr_text
                    
                    text += f"\n--- End of Page {page_num + 1} ---\n" #TODO: think about how to handle page ending in the future
                except Exception as page_error:
                    logging.error(f"Error processing page {page_num +1}: {page_error}")
                    text += f"\n[Error extracting page {page_num + 1}]\n"

            return text

    except Exception as e:
        logging.error(f"Error processing PDF {pdf_file_path}: {e}")
        return None
    
def extract_text_from_docx(docx_file_path): #TODO: Are we missing other elements if only considering paragraphs?
    """
    Extract text from a .docx file.

    Args:
        docx_file_path (str): The path to the .docx file.

    Returns:
        str: The extracted text, or None if it fails.
    """
    try:
        logging.info(f"Opening DOCX file: {docx_file_path}")
        doc = Document(docx_file_path)
        paragraphs = []

        # Loop through each paragraph
        for paragraph_num, paragraph in enumerate(doc.paragraphs, start=1):
            # Check if the paragraph text contains unexpected formatting or artifacts
            if paragraph.text.strip():
                clean_text = paragraph.text.strip()
                #TODO: python has library to take out markup instruction, check this up
                if "<" in clean_text or ">" in clean_text: # Check for HTML or XML or markup
                    logging.warning(f"Potential formatting instructions detected in paragraph {paragraph_num}: {clean_text}")
                paragraphs.append(clean_text)
            else:
                logging.debug(f"Paragraph {paragraph_num} is empty and was skipped.")

        extracted_text = "\n".join(paragraphs)
        logging.info(f"Successfully extracted text from {len(paragraphs)} paragraphs.")
        return extracted_text

    except PermissionError as e:
        logging.warning(f"Permission denied for {docx_file_path}: {e}")
        return None
    
    except Exception as e:
        logging.error(f"Error extracting text from DOCX file {docx_file_path}: {e}")
        return None

def extract_text_from_doc_with_tika(doc_file_path):
    """
    Extracts text from a .doc file using Apache Tika server.

    Args:
        doc_file_path (str): Path to the .doc file.

    Returns:
        str: Extracted text, or None if extraction fails.
    """
    try:
        parsed = parser.from_file(doc_file_path)  # Send the file to Tika server
        return parsed.get("content")  # Extract the content field (text)
    except Exception as e:
        print(f"Error extracting text from {doc_file_path}: {e}")
        return None

# def convert_doc_to_docx(doc_file_path):
#     """
#     Converts a .doc file to .docx format using the doc2docx library.

#     Args:
#         doc_file_path (str): Path to the .doc file to be converted.

#     Returns:
#         str: Path to the converted .docx file, or None if conversion fails.
#     """
#     try:
#         # Define the output file name
#         base_file_name = os.path.splitext(os.path.basename(doc_file_path))[0]
#         newdocx_file_path = os.path.join(os.path.dirname(doc_file_path), f"{base_file_name}.docx")

#         # Convert the file using doc2docx
#         logging.info(f"Converting {doc_file_path} to {newdocx_file_path}")
#         result = convert(doc_file_path, newdocx_file_path)

#         # Add this function to handle doc that has permission errors
#         if result.get('result') == 'error':
#             error_message = result.get('error', 'Unknown error')    # Add a fall back mechanism
#             logging.error(f"Conversion failed for {doc_file_path}: {error_message}")
#             return None
            
#         if os.path.exists(newdocx_file_path):
#             logging.info(f"Successfully converted {doc_file_path} to {newdocx_file_path}")
#             return newdocx_file_path
#         else:
#             logging.error(f"Conversion failed. Converted file not found at {newdocx_file_path}")
#             return None

#     except Exception as e:
#         logging.error(f"Error converting {doc_file_path} to .docx: {e}")
#         return None

# def extract_text_from_doc(doc_file_path):
#     """
#     Extract text from a .doc file. It will be converted to .docx before extraction.

#     Args:
#         doc_file_path (str): Path to the .doc file.

#     Returns:
#         str: Extracted text, or None if extraction fails.
#     """
#     try:
#         # Explicitly grant read and write permissions before conversion
#         try:
#             os.chmod(doc_file_path, stat.S_IRUSR | stat.S_IWUSR)
#             logging.info(f"Granted read and write permissions for {doc_file_path}")
#         except PermissionError as pe:
#             logging.error(f"Failed to grant permissions for {doc_file_path}: {pe}")
#             return None

#         # Convert .doc to .docx
#         if doc_file_path.endswith('.doc'):
#             newdocx_file_path = convert_doc_to_docx(doc_file_path)
#             if not newdocx_file_path:
#                 logging.error(f"Conversion of {doc_file_path} failed.")
#                 return None
            
#             # Delete the original .doc file after successful conversion
#             os.remove(doc_file_path)
#             logging.info(f"Deleted original .doc file: {doc_file_path}")

#             doc_file_path = newdocx_file_path # Use the converted .docx file for extraction

#         doc = Document(doc_file_path)
#         paragraphs = []

#         # Skip empty paragraph
#         for para_num, para in enumerate(doc.paragraphs, start=1):
#             if para.text.strip():
#                 paragraphs.append(para.text.strip())
#             else:
#                 logging.debug(f"Skipped empty paragraph {para_num}.")
            
#         extracted_text = "\n".join(paragraphs)
#         logging.info(f"Successfully extracted {len(paragraphs)} paragraphs from {doc_file_path}")

#         # Clean up: delete the .docx file after extraction
#         os.remove(doc_file_path)
#         logging.info(f"Deleted converted .docx file: {doc_file_path}")

#         return extracted_text
    
#     except Exception as e:
#         logging.error(f"Error extracting text from {doc_file_path}: {e}")
#         return None

def extract_text_from_txt(txt_file_path):
    """
    Extract text from a plain txt file using Python's built-in open() function.

    Args:
        txt_file_path (str): The path to the txt file.
    
    Returns:
        str: The extracted text, or None if it fails.
    """
    try:
        with open(txt_file_path, 'r', encoding='utf-8') as file:
            text = file.read()
        
        if not text.strip():
            logging.warning(f"Warning: The file {txt_file_path} is empty.")

        return text
    except Exception as e:
        logging.error(f"Error extracting text from TXT file: {e}")
        return None

def extract_text_from_pptx(pptx_file_path): #TODO: .ppt not supported, check tika and how to integrate 
    """
    Extract text from a PowerPoint file, including slide text, speaker notes, and tables.

    Args:
        pptx_file_path (str): The path to the PowerPoint file.

    Returns:
        str: The extracted text, or None if it fails.
    """
    try:
        logging.info(f"Opening PowerPoint file: {pptx_file_path}")
        presentation = pptx.Presentation(pptx_file_path)
        all_slides_text = []

        for slide_num, slide in enumerate (presentation.slides, start=1):
            slide_text = []

            for shape in slide.shapes:
                # Extract text from text-containing shapes
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # Extract text from table
                if shape.shape_type == 19: # Type 19 matches to a table
                    for row in shape.table.rows:
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                slide_text.append(cell_text)
            
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_text.append(f"Speaker Notes: {notes}")

            logging.info(f"Processed slide {slide_num}/{len(presentation.slides)}")
            
            if slide_text:
                # Combien slide content with headers
                all_slides_text.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
            else:
                logging.warning(f"Slide {slide_num} contains no extractable text.")
            
        return "\n\n".join(all_slides_text)
    
    except Exception as e:
        logging.error(f"Error extracting text from PowerPoint file {pptx_file_path}: {e}")
        return None

def process_files(api_url, endpoint, api_key, output_csv, download_path):
    """
    Access files from an API, download them, extract text, and log results.
    
    Args:
        api_url (str): Base URL for the API.
        endpoint (str): API endpoint for fetching records.
        api_key (str): My personal authentication key for API requests, please switch to yours.
        output_csv (str): Path to the output CSV.
        download_path (str): Directory to save downloaded files.
    """
    try:
        auth_headers = {'Authorization': f'Bearer {api_key}'}
        page = 1
        page_size = 100 # Invenio limit
        has_more_pages = True
        counter = 0

        # Updated headers to include 'DOI' and 'Flag' in the main CSV
        output_headers = ['Record ID', 'DOI', 'Languages', 'File Name', 'Extracted Text', 'Flag'] #TODO: change column name flag to 'failed' and add column 'supported'

        # Make sure output directories and CSV exists.
        os.makedirs(download_path, exist_ok=True)

        if not os.path.exists(output_csv):
            with open(output_csv, 'w', newline='', encoding='utf-8') as file:
                csv.writer(file).writerow(output_headers)

        while has_more_pages:
            response = requests.get(f"{api_url}/{endpoint}?size={page_size}&page={page}", headers=auth_headers)
            response.raise_for_status() # Raise an error for non-200 responses
            data = response.json()

            #Check each page and all its records
            for record in data.get('hits', {}).get('hits', []):
                record_id = record['id']
                doi = record.get('pids', {}).get('doi', {}).get('identifier', 'N/A')  # Extract DOI if available
                languages = ','.join([lang['id'] for lang in record['metadata'].get('languages', [])])
                
                logging.info(f"Processing record {record_id} (DOI: {doi}, Languages: {languages})")

                for file_name, file_metadata in record.get('files', {}).get('entries', {}).items():
                    counter += 1
                    logging.info(f"Processing file Number {counter}: {file_name}")

                    file_url = f"{api_url}/api/records/{record_id}/files/{file_name}/content"
                    local_file_path = os.path.join(download_path, file_name)

                    # Download the file
                    download_result = download_file(file_url, local_file_path, auth_headers)
                    if not download_result:
                        logging.error(f"Failed to download file {file_name} for record {record_id}.")
                        with open(output_csv, 'a', newline='', encoding='utf-8') as file:
                            writer = csv.writer(file)
                            writer.writerow ([record_id, doi, languages, file_name, '[Download Failed]', 1])
                        continue
        
                    # Extract file content and get flag
                    extracted_text, flag = extract_file(local_file_path, file_name)

                    # Append results to the output CSV
                    with open(output_csv, 'a', newline='', encoding='utf-8') as file:
                        writer = csv.writer(file)
                        writer.writerow([record_id, doi, languages, file_name, extracted_text or '[Error: Extraction Failed]', flag])
                        if flag == 0:
                            # Successfully processed, delete the file
                            os.remove(local_file_path)
                            logging.info(f"Deleted successfully processed file: {file_name}")
                        else:
                            # Retain files that could not be processed
                            logging.warning(f"Failed to process file {file_name} for record {record_id}. Retained in folder.")

            # Check if there are more pages to fetch
            logging.info(f"Completed page {page}.")
            has_more_pages = 'next' in data.get('links', {})
            page += 1

        logging.info("All files processed and results written to CSV.")
        #TODO: this logging info won't get run. If there is an error on one page, it can still go ot the next page.

    except requests.exceptions.HTTPError as e:
        logging.error(f"HTTP error: {e}")
    except Exception as e:
        logging.error(f"An error occurred: {e}")

def main():
    """Main function that leads the workflow"""
    api_url = "https://works.hcommons.org"
    api_key = "my api"
    api_endpoint = "api/records"
    output_csv = "output9.csv"
    download_path = "download_files9"

    os.makedirs(download_path, exist_ok=True)
    process_files(api_url, api_endpoint, api_key, output_csv, download_path)

if __name__ == "__main__":
    main()